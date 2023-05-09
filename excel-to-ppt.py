import os
import tempfile
from PIL import ImageGrab
import openpyxl
from pptx import Presentation
from pptx.util import Inches

def excel_sheet_to_image(sheet, output_file):
    # Save the current screen scaling settings
    screen_scaling = ImageGrab.screen_scaling()
    
    # Set the screen scaling to 100%
    ImageGrab.set_screen_scaling(1)
    
    # Grab the image of the sheet
    img = ImageGrab.grabclipboard(sheet)

    # Save the image to the output file
    img.save(output_file)

    # Reset the screen scaling
    ImageGrab.set_screen_scaling(screen_scaling)

def excel_to_ppt(input_excel, output_ppt):
    # Load the input Excel workbook
    wb = openpyxl.load_workbook(input_excel, read_only=True)

    # Create a new PowerPoint presentation
    ppt = Presentation()

    # Loop through each sheet in the Excel workbook
    for sheet in wb:
        # Create a new slide in the PowerPoint presentation
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])

        # Create a temporary file for the sheet image
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            # Convert the sheet to an image and save it to the temporary file
            excel_sheet_to_image(sheet, tmp_img.name)

            # Add the image to the PowerPoint slide
            slide.shapes.add_picture(tmp_img.name, Inches(0), Inches(0), width=ppt.slide_width, height=ppt.slide_height)

            # Remove the temporary file
            os.unlink(tmp_img.name)

    # Save the PowerPoint presentation
    ppt.save(output_ppt)

if __name__ == "__main__":
    input_excel = "your_input_excel_file.xlsx"
    output_ppt = "your_output_ppt_file.pptx"
    excel_to_ppt(input_excel, output_ppt)
